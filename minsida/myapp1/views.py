from django.contrib.auth import authenticate, login
from django.contrib.auth import logout
from django.contrib.auth.decorators import login_required
from django.shortcuts import render, redirect
from .forms import pu, RedigeraForm, Update
from .models import pulistan
from openpyxl import load_workbook
import subprocess

def home (request):
    if request.user.is_authenticated :
        return render(request, "test.html")
    else:
        return render(request, "login.html")

def psfiber (request):
    if request.user.is_authenticated :
        context = pulistan.objects.all()
        return render(request, 'psfiber.html', {'context' : context})

    else:
        return render(request, "login.html")

def nyttpu(request) :
    form = pu
    if request.method == 'POST':
        form = pu(request.POST)
        if form.is_valid() :
            bla = pulistan(
                avd=form.cleaned_data["avd"],
                beskrivning =form.cleaned_data["beskrivning"],
                artal=form.cleaned_data["artal"],
                kostnad=form.cleaned_data["kostnad"],
                projektnr=form.cleaned_data["projektnr"],
            )
            bla.save()
        return redirect('psfiber')
    return render (request, 'nyttpu.html', {'form': form})

def redigera(request, item_id):
    bla = pulistan.objects.get(pk=item_id)
    avd = bla.avd
    beskrivning = bla.beskrivning
    artal = bla.artal
    kostnad = bla.kostnad
    projektnr = bla.projektnr
    formular = RedigeraForm(initial={'avd': avd,'beskrivning': beskrivning, 'kostnad' : kostnad,'artal': artal, 'projektnr': projektnr})
    context = pulistan.objects.filter(id=item_id).distinct()
    if request.method == 'POST':
        form = Update(request.POST, instance=bla)
        if form.is_valid() :
            form.save()
        return redirect('psfiber')
    return render(request, 'redigera.html', {'form' : formular, 'context' : context})

def remove(request, item_id) :
    bla.pulistan.objects.get(pk=item_id)
    item.delete()
    previous_page = request.META['HTTP_REFERER']
    return redirect(previous_page)

def skapa_askande(request, item_id) :
    bla = pulistan.objects.get(pk=item_id)

    wb = load_workbook('test1/resources/pu_mall1.xlsx')

    # grab the active worksheetcd
    ws = wb.active

    # Data can be assigned directly to cells
    ws['B6'] = bla.avd
    ws['D6'] = bla.beskrivning
    ws['J6'] = bla.projektnr
    ws['J6'] = bla.projektnr

    # Save the file
    filnamn = bla.beskrivning
    wb.save(f'test1/resources/{filnamn}.xlsx')
    subprocess.check_call(['open', '-a', 'Microsoft Excel', f'test1/resources/{filnamn}.xlsx'])
    bla.askande_fil = (f'test1/resources/{filnamn}.xlsx')
    bla.save()
    previous_page = request.META['HTTP_REFERER']
    return redirect(previous_page)


def user_login(request):
    if request.method == 'POST':
        # Process the request if posted data are available
        username = request.POST['username']
        password = request.POST['password']
        # Check username and password combination if correct
        user = authenticate(username=username, password=password)
        if user is not None:
            # Save session as cookie to login the user
            login(request, user)
            # Success, now let's login the user.
            return render(request, 'psfiber.html')
        else:
            # Incorrect credentials, let's throw an error to the screen.
            return render(request, 'login.html', {'error_message': 'Incorrect username and / or password.'})
    else:
        # No post data availabe, let's just show the page to the user.
        return render(request, 'login.html')

def user_logout(request):
    logout(request)
    return render(request, 'login.html')

def vrapport(request):
    import pandas as pd
    import numpy as np
    import tempfile
    import csv
    from numerize import numerize

    if request.method == 'POST':
     mf = request.FILES.get('myfile')
     df = pd.read_csv(mf,sep=";", error_bad_lines=False, header=0)
     rubrik = int(len(df.index))
     df['Planerad kostnad'] = df['Planerad kostnad'].str.replace(" ", "",).str.replace(",", ".").astype(float)
     kostnad = numerize.numerize(df['Planerad kostnad'].sum(), 1) + "SEK"
     avd13_kostn =  numerize.numerize(df.loc[df['Avdelning'] == 13, 'Planerad kostnad'].sum(), 1) + "SEK"
     avd13_ao = len(df.loc[df['Avdelning'] == 13])
     avd14_kostn = numerize.numerize(df.loc[df['Avdelning'] == 14, 'Planerad kostnad'].sum(), 1) + "SEK"
     avd14_ao = len(df.loc[df['Avdelning'] == 14])
     avd15_kostn = numerize.numerize(df.loc[df['Avdelning'] == 15, 'Planerad kostnad'].sum(), 1) + "SEK"
     avd15_ao = len(df.loc[df['Avdelning'] == 15])
     avd18_kostn = numerize.numerize(df.loc[df['Avdelning'] == 18, 'Planerad kostnad'].sum(), 1) + "SEK"
     avd18_ao = len(df.loc[df['Avdelning'] == 18])

     rk0 = len(df[df['Riskkod'] == 'RK0'])
     rk1 = len(df[df['Riskkod'] == 'RK1'])
     rk2 = len(df[df['Riskkod'] == 'RK2'])
     rk3 = len(df[df['Riskkod'] == 'RK3'])
     rk4 = len(df[df['Riskkod'] == 'RK4'])
     df['Reg.datum'] = pd.to_datetime(df['Reg.datum'])
     jun = len(df[df['Reg.datum'].dt.month <= 6])
     jul = len(df[df['Reg.datum'].dt.month <= 7])
     aug = len(df[df['Reg.datum'].dt.month <= 8])
     sep = len(df[df['Reg.datum'].dt.month <= 9])
     okt = len(df[df['Reg.datum'].dt.month <= 10])
     nov = len(df[df['Reg.datum'].dt.month <= 11])
     dec = len(df[df['Reg.datum'].dt.month <= 12])
     jan = len(df[df['Reg.datum'].dt.month <= 1])


     return render(request, 'vrapport.html', {
         'rubrik' : rubrik,
         'kostnad' : kostnad,
         'rk0' : rk0,
         'rk1' : rk1,
         'rk2' : rk2,
         'rk3' : rk3,
         'rk4' : rk4,
         'jun' : jun,
         'jul' : jul,
         'aug' : aug,
         'sep' : sep,
         'okt' : okt,
         'nov' : nov,
         'dec' : dec,
         'jan' : jan,
         'avd13_kostn' : avd13_kostn,
         'avd13_ao' : avd13_ao,
         'avd14_kostn' : avd14_kostn,
         'avd14_ao' : avd14_ao,
         'avd15_kostn' : avd15_kostn,
         'avd15_ao' : avd15_ao,
         'avd18_kostn' : avd18_kostn,
         'avd18_ao' : avd18_ao,
     })
    else:

      return render(request, 'vrapport.html')

def skaparapport(request):
    from pptx import Presentation


    wb_obj = load_workbook('test1/resources/RCA.xlsx')

    sheet_obj = wb_obj.active

    cell_obj = sheet_obj.cell(row=4, column=7)

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = cell_obj.value
    subtitle.text = "python-pptx was here!"
    prs.save(f'test1/resources/test.pptx')

    return render(request, 'vrapport.html')